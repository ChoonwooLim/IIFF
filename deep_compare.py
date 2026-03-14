"""
Deep comparison between PPTX and HTML content.
Extracts all text from both, does slide-by-slide mapping, and finds all differences.
"""
from pptx import Presentation
from html.parser import HTMLParser
import re

# ===== EXTRACT PPTX TEXT =====
pptx_path = r'D:\GoogleDrive\인천국제영화제\20260314_ iiff기획 정리.pptx'
prs = Presentation(pptx_path)

pptx_slides = []
for i, slide in enumerate(prs.slides, 1):
    slide_texts = []
    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.text.strip():
            slide_texts.append(shape.text.strip())
        if shape.has_table:
            tbl = shape.table
            rows = []
            for row in tbl.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(' | '.join(cells))
            slide_texts.append('[TABLE]\n' + '\n'.join(rows))
    pptx_slides.append({
        'num': i,
        'texts': slide_texts,
        'full_text': '\n'.join(slide_texts)
    })

# ===== EXTRACT HTML SLIDE TEXT =====
class SlideExtractor(HTMLParser):
    def __init__(self):
        super().__init__()
        self.slides = []
        self.current_slide = None
        self.current_text = []
        self.in_script = False
        self.in_style = False
        self.in_nav = False
        self.depth = 0
        
    def handle_starttag(self, tag, attrs):
        attrs_dict = dict(attrs)
        classes = attrs_dict.get('class', '')
        
        if tag == 'script':
            self.in_script = True
        if tag == 'style':
            self.in_style = True
        if 'nav-bar' in classes:
            self.in_nav = True
            
        if tag == 'div' and 'slide' in classes.split():
            if self.current_slide is not None:
                self.slides.append({
                    'classes': self.current_slide,
                    'text': ' '.join(self.current_text).strip()
                })
            self.current_slide = classes
            self.current_text = []
    
    def handle_endtag(self, tag):
        if tag == 'script':
            self.in_script = False
        if tag == 'style':
            self.in_style = False
            
    def handle_data(self, data):
        if self.in_script or self.in_style or self.in_nav:
            return
        if self.current_slide is not None:
            text = data.strip()
            if text:
                self.current_text.append(text)
    
    def finalize(self):
        if self.current_slide is not None:
            self.slides.append({
                'classes': self.current_slide,
                'text': ' '.join(self.current_text).strip()
            })

with open(r'C:\WORK\IIFF\presentation_full.html', 'r', encoding='utf-8') as f:
    html_content = f.read()

parser = SlideExtractor()
parser.feed(html_content)
parser.finalize()

print("=" * 80)
print(f"PPTX: {len(pptx_slides)} slides")
print(f"HTML: {len(parser.slides)} slides")
print("=" * 80)

# Print PPTX slide summaries
print("\n\n===== PPTX SLIDE STRUCTURE =====")
for s in pptx_slides:
    first_line = s['full_text'][:120].replace('\n', ' | ') if s['full_text'] else '(empty/image only)'
    print(f"  PPTX Slide {s['num']:2d}: {first_line}")

print("\n\n===== HTML SLIDE STRUCTURE =====")
for i, s in enumerate(parser.slides, 1):
    first_line = s['text'][:120] if s['text'] else '(empty)'
    print(f"  HTML Slide {i:2d} [{s['classes'][:40]}]: {first_line}")

# Detailed Personnel comparison
print("\n\n===== PERSONNEL COMPARISON =====")
print("\n--- PPTX Personnel (Slides 71-74) ---")
for s in pptx_slides[70:74]:
    print(f"\nSlide {s['num']}:")
    print(s['full_text'][:500])

print("\n--- HTML Personnel (last content slides) ---")
for i, s in enumerate(parser.slides):
    if 'Personnel' in s['text'] or '컨트롤 타워' in s['text'] or '글로벌 네트워크' in s['text'] or '실무 운영' in s['text'] or '대외 협력' in s['text']:
        print(f"\nHTML Slide {i+1}:")
        # Print names found
        names = re.findall(r'[가-힣]{2,4}\s*\(', s['text'])
        for n in names:
            print(f"  Found: {n}")

# Now let's find text that's in PPTX but NOT in HTML (new content)
print("\n\n===== NEW CONTENT IN PPTX (not in HTML) =====")
html_all_text = ' '.join(s['text'] for s in parser.slides)

new_phrases = [
    "if는 일어날수 있는 모든 우연 즉 모든 가능성입니다",
    "끝없는 상상입니다",
    "헐리우드와 공동제작, 투자, 배급 실전협업 영화제",
    "상업영화와 독립영화가 공존하는 영화제",
    "참여하고 공동창작하고 즐기는 잼보리형 영화제",
    "온라인 오프라인 365일",
    "영화 · 음악 · 테크놀로지 · K-컬처가 융합된 영화 플랫폼",  # already in overview slide
    "Cinematic Insights from INSPIRE Resort",
    "See, enjoy, create, evaluate each other",
    "We meet across borders",
    "The one and only film festival in the universe",
    "Incheon will now become the gateway",
    "BIFF 노하우를 가진",
    "이청산",
    "박병용",
    "오석근",
    "송승희",
    "비젼",  # PPTX uses 비젼 instead of 비전
    "실현 전략",  # PPTX uses 실현 instead of 구현
    "All coincidences, all possibilities",
    "Unlimited imagination and free creation",
    "complex platform",
    "Together Platform",
    "The center of K - content",
    "Amazing Inspire Resort",
    "Attractive K-pop, K-food",
    "Mobile Short Film Competition + Camping Festival Viral Engine",
    "365 days a year, offline and online",  
    "networking - IIFF Partner Circle",
    "long-term partnership",
    "실행전략",
]

for phrase in new_phrases:
    if phrase in html_all_text:
        print(f"  ✅ FOUND in HTML: '{phrase}'")
    else:
        print(f"  ❌ MISSING from HTML: '{phrase}'")

# Check PPTX subtitles vs HTML
print("\n\n===== PPTX SUBTITLES (English) vs HTML =====")
for s in pptx_slides:
    for t in s['texts']:
        if t.startswith('[PLACEHOLDER'):
            # Extract subtitle text
            match = re.search(r'\] (.+)', t)
            if match:
                subtitle = match.group(1).strip()
                if subtitle and len(subtitle) > 5:
                    if subtitle in html_all_text:
                        pass  # found
                    else:
                        print(f"  Slide {s['num']}: MISSING subtitle: '{subtitle[:80]}'")
