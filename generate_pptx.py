"""
Hybrid HTML-to-PPTX generator.
Uses Playwright to capture high-res backgrounds without text,
then overlays native editable TextBoxes at exact coordinates using python-pptx.
"""
import sys, io, json, os, re
from playwright.sync_api import sync_playwright
import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

presentation_path = r'file:///C:/WORK/IIFF/presentation_full.html'
out_dir = r'C:/WORK/IIFF/tmp_slides'
out_pptx = r'C:/WORK/IIFF/IIFF_NextWave_Hybrid.pptx'

os.makedirs(out_dir, exist_ok=True)

# Helper to parse rgba/rgb to tuple
def parse_color(color_str):
    if not color_str: return (0, 0, 0)
    match = re.search(r'rgba?\((\d+),\s*(\d+),\s*(\d+)', color_str)
    if match:
        return (int(match.group(1)), int(match.group(2)), int(match.group(3)))
    return (0, 0, 0) # Fallback to black

def generate_hybrid_pptx():
    print("Starting Playwright to extract layout...")
    
    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        # 16:9 1920x1080 viewport for high quality
        page = browser.new_page(viewport={"width": 1920, "height": 1080})
        page.goto(presentation_path, wait_until="networkidle")
        
        # Inject script to hide scrollbars and prepare slider
        page.evaluate("""
            document.body.style.overflow = 'hidden';
            var slider = document.getElementById('slider');
            if(slider) {
                slider.style.width = '1920px';
                slider.style.height = '1080px';
                slider.style.position = 'absolute';
            }
        """)
        
        slides_count = page.evaluate("document.querySelectorAll('.slide').length")
        print(f"Detected {slides_count} slides.")
        
        slide_data = [] # List of dicts: {'bg_path': str, 'texts': list}
        
        for i in range(slides_count):
            print(f"Processing Slide {i+1}...")
            
            # Switch to this slide exactly
            page.evaluate(f"""
                () => {{
                    var slides = document.querySelectorAll('.slide');
                    slides.forEach((s, idx) => {{
                        if(idx === {i}) {{
                            s.classList.add('active');
                            s.style.opacity = '1';
                            s.style.visibility = 'visible';
                            s.style.zIndex = '100';
                            s.style.display = 'flex';
                        }} else {{
                            s.classList.remove('active');
                            s.style.opacity = '0';
                            s.style.visibility = 'hidden';
                            s.style.zIndex = '0';
                            s.style.display = 'none';
                        }}
                    }});
                }}
            """)
            
            # Wait a tiny bit for transition/render
            page.wait_for_timeout(300)
            
            # Extract text elements
            texts_json = page.evaluate(f"""
                () => {{
                    const slide = document.querySelectorAll('.slide')[{i}];
                    const walker = document.createTreeWalker(slide, NodeFilter.SHOW_TEXT, null, false);
                    let node;
                    let results = [];
                    
                    // Specific class to exclude
                    const excludeClasses = ['slide-number']; 
                    
                    while (node = walker.nextNode()) {{
                        if (!node.nodeValue.trim()) continue;
                        const el = node.parentElement;
                        if (!el || el.tagName === 'SCRIPT' || el.tagName === 'STYLE') continue;
                        
                        let isExcluded = false;
                        excludeClasses.forEach(c => {{ if(el.classList.contains(c)) isExcluded = true; }});
                        if(isExcluded) continue;
                        
                        const rect = el.getBoundingClientRect();
                        const style = window.getComputedStyle(el);
                        
                        // Ignore extremely hidden things
                        if (style.opacity === '0' || style.visibility === 'hidden' || style.display === 'none') continue;
                        if (rect.width === 0 || rect.height === 0) continue;
                        if (rect.top > 1500) continue; // Offscreen
                        
                        results.push({{
                            text: node.nodeValue.trim(),
                            x: rect.left,
                            y: rect.top,
                            w: rect.width,
                            h: rect.height,
                            color: style.color,
                            fontSize: parseFloat(style.fontSize),
                            fontWeight: style.fontWeight,
                            textAlign: style.textAlign,
                            elementId: Math.random().toString(36).substr(2, 9) // temporary ID
                        }});
                        
                        el.setAttribute('data-temp-id', results[results.length-1].elementId);
                    }}
                    return results;
                }}
            """)
            
            # Make text transparent
            page.evaluate("""
                document.querySelectorAll('[data-temp-id]').forEach(el => {
                    el.style.color = 'transparent';
                    // also hide list bullets if needed, but transparent color usually covers it
                });
            """)
            page.wait_for_timeout(100) # wait for DOM update
            
            # Take background screenshot
            bg_path = os.path.join(out_dir, f'slide_bg_{i+1}.png')
            page.screenshot(path=bg_path)
            
            # Restore text color
            page.evaluate("""
                document.querySelectorAll('[data-temp-id]').forEach(el => {
                    el.style.color = '';
                    el.removeAttribute('data-temp-id');
                });
            """)
            
            slide_data.append({
                'bg_path': bg_path,
                'texts': texts_json
            })
            
        browser.close()
        
    print("Playwright extraction complete. Building PPTX...")
    
    # PPTX Generation
    prs = pptx.Presentation()
    # 16:9 ratio. 16 inches x 9 inches
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    blank_slide_layout = prs.slide_layouts[6] # Blank layout
    
    for idx, sdata in enumerate(slide_data):
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Add background image
        slide.shapes.add_picture(sdata['bg_path'], 0, 0, width=Inches(16), height=Inches(9))
        
        # Add text overlays
        for t in sdata['texts']:
            # Screen was 1920x1080. 16 inches x 9 inches.
            # Convert px to Inches -> (pixels / 1920) * 16
            left = Inches((t['x'] / 1920.0) * 16.0)
            top = Inches((t['y'] / 1080.0) * 9.0)
            
            # Add safe padding to width/height to prevent text wrapping in PPT
            width = Inches(((t['w'] + 50) / 1920.0) * 16.0)
            height = Inches(((t['h'] + 20) / 1080.0) * 9.0)
            
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = False
            
            # Add paragraph
            p = tf.paragraphs[0]
            p.text = t['text']
            
            # Determine font pt. 1080px = 9 inches = 648 pt. => 1px = 648/1080 pt = 0.6pt
            font_pt = t['fontSize'] * 0.6
            p.font.size = Pt(font_pt)
            p.font.name = 'Inter' # Default to modern font
            
            # Color
            r, g, b = parse_color(t['color'])
            p.font.color.rgb = RGBColor(r, g, b)
            
            # Bold
            if str(t['fontWeight']) in ['bold', '600', '700', '800', '900']:
                p.font.bold = True
                
            # Alignment
            if t['textAlign'] == 'center':
                p.alignment = PP_ALIGN.CENTER
            elif t['textAlign'] == 'right':
                p.alignment = PP_ALIGN.RIGHT
            else:
                p.alignment = PP_ALIGN.LEFT
                
            # PPTX TextBox default margins can shift text visually. Reset them.
            txBox.margin_left = 0
            txBox.margin_right = 0
            txBox.margin_top = 0
            txBox.margin_bottom = 0
            
    prs.save(out_pptx)
    print(f"Successfully saved {slides_count} slides to {out_pptx}")

if __name__ == '__main__':
    generate_hybrid_pptx()
