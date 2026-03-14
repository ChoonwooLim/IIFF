from pptx import Presentation
from pptx.util import Inches, Pt
import os

pptx_path = r'D:\GoogleDrive\인천국제영화제\20260314_ iiff기획 정리.pptx'
prs = Presentation(pptx_path)

output = []
for i, slide in enumerate(prs.slides, 1):
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.text.strip():
            texts.append(f"  [{shape.shape_type}] {shape.text.strip()[:300]}")
        if shape.has_table:
            tbl = shape.table
            for row_idx, row in enumerate(tbl.rows):
                row_texts = []
                for cell in row.cells:
                    row_texts.append(cell.text.strip())
                texts.append(f"  [TABLE ROW {row_idx}] {' | '.join(row_texts)}")
    
    slide_text = '\n'.join(texts) if texts else '  (no text found)'
    output.append(f"===== SLIDE {i} =====\n{slide_text}")

out_path = r'C:\WORK\IIFF\pptx_full_content.txt'
with open(out_path, 'w', encoding='utf-8') as f:
    f.write('\n\n'.join(output))

print(f"Extracted {len(prs.slides)} slides -> {out_path}")
