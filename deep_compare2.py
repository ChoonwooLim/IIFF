"""Deep comparison between PPTX and HTML content."""
import sys
import io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from html.parser import HTMLParser
import re

# ===== EXTRACT PPTX =====
prs = Presentation(r'D:\GoogleDrive\인천국제영화제\20260314_ iiff기획 정리.pptx')

pptx_slides = []
for i, slide in enumerate(prs.slides, 1):
    slide_texts = []
    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.text.strip():
            slide_texts.append(shape.text.strip())
        if shape.has_table:
            tbl = shape.table
            rows = []
            for row in tbl.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(' | '.join(cells))
            slide_texts.append('[TABLE]\n' + '\n'.join(rows))
    pptx_slides.append({'num': i, 'texts': slide_texts, 'full': '\n'.join(slide_texts)})

# ===== EXTRACT HTML =====
class SlideExtractor(HTMLParser):
    def __init__(self):
        super().__init__()
        self.slides = []
        self.cur = None
        self.texts = []
        self.skip = False
    def handle_starttag(self, tag, attrs):
        d = dict(attrs)
        c = d.get('class', '')
        if tag in ('script', 'style'): self.skip = True
        if 'nav-bar' in c: self.skip = True
        if tag == 'div' and 'slide' in c.split():
            if self.cur is not None:
                self.slides.append({'cls': self.cur, 'text': ' '.join(self.texts).strip()})
            self.cur = c; self.texts = []
    def handle_endtag(self, tag):
        if tag in ('script', 'style'): self.skip = False
    def handle_data(self, data):
        if not self.skip and self.cur is not None:
            t = data.strip()
            if t: self.texts.append(t)
    def done(self):
        if self.cur: self.slides.append({'cls': self.cur, 'text': ' '.join(self.texts).strip()})

with open(r'C:\WORK\IIFF\presentation_full.html', 'r', encoding='utf-8') as f:
    html = f.read()
p = SlideExtractor(); p.feed(html); p.done()

print(f"PPTX: {len(pptx_slides)} slides, HTML: {len(p.slides)} slides\n")

# ===== CONTENT CHECKS =====
html_all = ' '.join(s['text'] for s in p.slides)

checks = [
    # (label, search_text)
    ("PPTX Slide 2: 헐리우드와 공동제작, 투자, 배급 실전협업 영화제", "헐리우드와 공동제작, 투자, 배급 실전협업"),
    ("PPTX Slide 2: 상업영화와 독립영화가 공존하는 영화제", "상업영화와 독립영화가 공존하는 영화제"),
    ("PPTX Slide 2: 참여하고 공동창작하고 즐기는 잼보리형", "참여하고 공동창작하고 즐기는 잼보리형"),
    ("PPTX Slide 2: 온라인 오프라인 365일 영화제", "온라인 오프라인 365일"),
    ("PPTX Slide 2: 영화·음악·테크놀로지·K-컬처가 융합된", "영화 · 음악 · 테크놀로지 · K-컬처가 융합된 영화 플랫폼"),
    ("PPTX Slide 2: Cinematic Insights from INSPIRE", "Cinematic Insights from INSPIRE"),
    ("PPTX Slide 4: if는 일어날수 있는 모든 우연", "if는 일어날수 있는 모든 우연"),
    ("PPTX Slide 4: 끝없는 상상입니다", "끝없는 상상입니다"),
    ("PPTX Slide 5: Unlimited imagination", "Unlimited imagination"),
    ("PPTX Slide 6: See, enjoy, create", "See, enjoy, create"),
    ("PPTX Slide 6: We meet across borders", "We meet across borders"),
    ("PPTX Slide 6: one and only film festival in the universe", "one and only film festival in the universe"),
    ("PPTX Slide 6: Incheon will now become the gateway", "Incheon will now become the gateway"),
    ("PPTX Slide 7: 00. complex platform", "complex platform"),
    ("PPTX Slide 12: Together Platform subtitle", "Together Platform"),
    ("PPTX Slide 13: The center of K - content", "The center of K"),
    ("PPTX Slide 14: Amazing Inspire Resort", "Amazing Inspire Resort"),
    ("PPTX Slide 15: Attractive K-pop, K-food", "Attractive K-pop"),
    ("PPTX Slide 16: Mobile Short Film Competition subtitle", "Mobile Short Film Competition + Camping"),
    ("PPTX Slide 17: 365 days a year subtitle", "365 days a year"),
    ("PPTX Slide 18: networking - IIFF Partner Circle", "networking - IIFF Partner Circle"),
    ("PPTX Slide 19: long-term partnership subtitle", "long-term partnership"),
    ("PPTX Slide 37: 실행전략 · strategy", "실행전략"),
    ("PPTX Slide 38: 영화제 비젼 (not 비전)", "비젼"),
    ("PPTX Slide 38: 구체적 실현 전략 (not 구현)", "구체적 실현 전략"),
    ("PPTX Slide 45: A-to-Z → 5. A-to-Z Roadmap subtitle", "전체 예산 30억 원 기준"),
    ("PPTX Slide 51: BIFF Comparison subtitle", "선배 영화제의 성과를"),
    ("PPTX Slide 66: Risk Mgmt subtitle", "先 민간 주도"),
    ("Person: 이청산", "이청산"),
    ("Person: 박병용", "박병용"),
    ("Person: 오석근", "오석근"),
    ("Person: 송승희", "송승희"),
    ("Person: 이용관", "이용관"),
    ("Person: 저스틴 김", "저스틴 김"),
    ("Person: 돈 플랑캔", "돈 플랑캔"),
    ("Person: 유동수", "유동수"),
    ("Person: 조광희", "조광희"),
    ("Person: 서태웅", "서태웅"),
    ("Person: 황보진호", "황보진호"),
    ("Person: 강준", "강준"),
    ("Person: 노준석", "노준석"),
    ("Person: 임춘우", "임춘우"),
    ("Person: 김무전", "김무전"),
    ("Person: 제니스", "제니스"),
]

print("===== CONTENT CHECK =====")
for label, search in checks:
    found = search in html_all
    mark = "OK" if found else "MISSING"
    print(f"  [{mark:7s}] {label}")

# Check PPTX sub-titles that are not in HTML
print("\n===== PPTX SUBTITLE ANALYSIS =====")
for s in pptx_slides:
    for t in s['texts']:
        # Subtitles in PPTX are in PLACEHOLDER elements
        if '[PLACEHOLDER' in str(t):
            continue
        lines = t.split('\n')
        for line in lines:
            line = line.strip()
            if line and len(line) > 10 and line not in html_all:
                # Check if it's a truncated match
                if line[:20] in html_all:
                    continue
                print(f"  Slide {s['num']}: NOT IN HTML: '{line[:100]}'")

# Check which PPTX slides have NO match in HTML
print("\n===== PPTX SLIDES WITH UNIQUE CONTENT =====")
for s in pptx_slides:
    if not s['full']:
        continue
    # Check first meaningful text line
    lines = [l.strip() for l in s['full'].split('\n') if l.strip() and len(l.strip()) > 5]
    if not lines:
        continue
    unmatched = []
    for line in lines[:3]:
        if line[:15] not in html_all:
            unmatched.append(line[:80])
    if unmatched:
        print(f"  Slide {s['num']}: {unmatched[0]}")
